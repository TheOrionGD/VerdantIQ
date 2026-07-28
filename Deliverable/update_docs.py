import pptx
import docx

def update_docx():
    docx_path = 'o:/PROJECTS/VerdantIQ/Deliverable/Write UP.docx'
    doc = docx.Document(docx_path)
    
    # Update Paragraph 7 (Project Title)
    for p in doc.paragraphs:
        if 'Title of the Project' in p.text:
            p.text = 'Title of the Project\t:\tEcoSphere (VerdantIQ)'
            
        if 'ABSTRACT' in p.text:
            continue
            
        if 'Conventional environmental sustainability platforms are limited' in p.text:
            p.text = (
                "     Conventional environmental sustainability platforms are limited by historical logging, manual data entry, and generic recommendations, leaving a critical research gap in active, personalized decision support and multi-resource optimization. To address these limitations, this project implements EcoSphere (VerdantIQ), a fully realized, machine learning-driven environmental decision support platform designed to transition sustainability management from passive tracking to proactive optimization.\n\n"
                "     The architecture employs a React 18 Spatial Canvas client interface equipped with the Emerald Oasis design system, a Groq AI floating Omnibar command center, and a 3-mode view switcher (Spatial Canvas, Executive Desktop Dashboard, and Touch Mobile App). The client communicates with a Spring Boot 3.2 core middleware gateway utilizing Spring Security with JWT authentication and Spring Data MongoDB. Spatial features are managed via MongoDB native 2dsphere GeoJSON indexing ($geoWithin campus geofencing verification), paired with MinIO / AWS S3 object storage integrated with Tesseract OCR for automated utility bill text parsing.\n\n"
                "     The intelligence tier consists of a Python 3.11 FastAPI microservice deploying XGBoost regression models for 30-day resource consumption forecasting, scikit-learn IsolationForest for usage anomaly detection, and a Google OR-Tools Mixed-Integer Linear Programming (MILP) solver to evaluate multivariable trade-offs between carbon emissions, utility usage, and financial budgets. Generative AI reasoning is powered by a Groq API Gateway (Llama-3.3-70b-versatile) equipped with a local grounded fallback engine to deliver explainable recommendations without service interruption. The platform also includes real-time Server-Sent Events (SSE) notification streaming, system telemetry monitoring, and an automated Apache PDFBox executive report generator. The complete system has been fully verified across a 36-endpoint automated integration test suite and orchestrated via a multi-service PowerShell launcher."
            )
            
        if 'KEYWORDS:' in p.text:
            p.text = "KEYWORDS: Environmental Decision Support, Machine Learning, Multi-Resource Optimization, Behavioral Forecasting, Digital Twin Simulation, Groq AI LLM Gateway, Google OR-Tools MILP, React Spatial Canvas, MongoDB 2dsphere, MinIO S3 OCR, Sustainability Intelligence"

    doc.save(docx_path)
    print(f"Successfully updated DOCX file: {docx_path}")


def update_pptx():
    pptx_path = 'o:/PROJECTS/VerdantIQ/Deliverable/EcoSphere_Proposal_Review.pptx'
    prs = pptx.Presentation(pptx_path)
    
    # Helper to replace text in shape paragraph while preserving font size if possible
    def replace_para_text(p, new_text):
        p.text = new_text

    # Slide 1: Title
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'EcoSphere' in p.text and 'Platform' in p.text:
                    p.text = 'EcoSphere (VerdantIQ) — Machine Learning-Driven Environmental Decision Support Platform'
                elif 'Proposed Project Review' in p.text or 'Implementation' in p.text:
                    p.text = 'Full System Architectural Implementation & Review'

    # Slide 2: Executive Summary
    s2 = prs.slides[1]
    for shape in s2.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Spring Boot gateway' in p.text or 'PostgreSQL' in p.text:
                    p.text = 'Architecture: React 18 Spatial Canvas shell, Spring Boot 3.2 REST gateway, MongoDB 2dsphere spatial database, MinIO S3 OCR storage, Python FastAPI ML container, Google OR-Tools MILP constraint engine, and Groq AI LLM Gateway.'

    # Slide 4: Research Gap
    s4 = prs.slides[3]
    for shape in s4.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Gemini/OpenAI' in p.text or 'PostgreSQL' in p.text:
                    p.text = 'Solution Positioning: VerdantIQ integrates React 18 Spatial Canvas, MongoDB 2dsphere spatial geofencing, Google OR-Tools MILP optimization, and Groq AI LLM decision explanations.'

    # Slide 5: Core Features
    s5 = prs.slides[4]
    for shape in s5.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Evidence Verification Framework' in p.text:
                    p.text = 'Evidence Verification & Proof Storage'
                elif 'Validates user green actions through geocoding' in p.text:
                    p.text = 'Audits utility bills via MinIO S3 presigned storage, Tesseract OCR text extraction, and MongoDB 2dsphere $geoWithin campus boundary verification.'

    # Slide 6: Workflow Pipeline
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'PostgreSQL' in p.text:
                    p.text = p.text.replace('PostgreSQL', 'MongoDB 2dsphere & MinIO S3')
                if 'Gemini' in p.text:
                    p.text = p.text.replace('Gemini', 'Groq AI Gateway')

    # Slide 7: Dual-Engine Architecture & 7-Layer Model
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Gemini/OpenAI' in p.text:
                    p.text = 'Layer 5: Cognitive Reasoning — Manages Groq API (Llama-3.3-70b) prompts & local grounded fallback for advice.'
                if 'geolocation, QR, and hash' in p.text:
                    p.text = 'Layer 7: Verification & Audit — Runs MongoDB 2dsphere geofencing, MinIO S3 proof storage, and Tesseract OCR.'

    # Slide 8: Technical Stack Table
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if shape.has_table:
            table = shape.table
            row_updates = [
                ("Frontend Framework", "React 18 Spatial Canvas + Emerald Oasis theme, Groq AI Omnibar, 3-Mode View Switcher."),
                ("Backend Gateway", "Spring Boot 3.2, Spring Security with JWT tokens, and Spring Data MongoDB."),
                ("Data & Storage", "MongoDB 7.0 (2dsphere spatial index), MinIO / AWS S3 object storage for bill proof photos and OCR."),
                ("Machine Learning", "FastAPI (Python 3.11), scikit-learn IsolationForest anomaly engine, XGBoost regression (30-day forecast)."),
                ("Optimization Solver", "Google OR-Tools Mixed-Integer Linear Programming (MILP solver with pywraplp)."),
                ("AI & External APIs", "Groq API Gateway (Llama-3.3-70b-versatile) + Local Grounded Fallback Engine, OpenWeather & AQICN APIs.")
            ]
            for r_idx, (layer_name, tech_desc) in enumerate(row_updates, start=1):
                if r_idx < len(table.rows):
                    table.rows[r_idx].cells[0].text = layer_name
                    table.rows[r_idx].cells[1].text = tech_desc
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'PostGIS Coordinates' in p.text:
                    p.text = 'MongoDB 2dsphere Geofences Campus Radius'
                if 'Gemini API Demystifies' in p.text:
                    p.text = 'Groq LLM Gateway Explains Decision Logic'

    # Slide 10: Interface & Design
    s10 = prs.slides[9]
    for shape in s10.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'React.js Frontend Client' in p.text:
                    p.text = 'React 18 Spatial Canvas Client'
                if 'Interactive Dashboards:' in p.text:
                    p.text = 'Spatial Canvas Workspace: 2D infinite workspace with drag-to-pan, SVG data streams, and 12 interactive nodes.'
                if 'AI Assistant Chat:' in p.text:
                    p.text = 'Groq AI Omnibar & Assistant: Floating command bar with natural language intent parser and instant chat.'

    # Slide 11: Backend Middleware
    s11 = prs.slides[10]
    for shape in s11.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'PostgreSQL databases' in p.text:
                    p.text = 'Data Persistence: Utilizes Spring Data MongoDB for document storage and 2dsphere GeoJSON spatial queries.'
                if 'Gemini API Connector' in p.text:
                    p.text = 'Groq LLM Gateway: Manages Groq API calls (Llama-3.3-70b) with local grounded fallback engine for offline reliability.'

    # Slide 12: Core Intelligent Capabilities Table
    s12 = prs.slides[11]
    for shape in s12.shapes:
        if shape.has_table:
            table = shape.table
            table.rows[4].cells[0].text = "Explainable AI"
            table.rows[4].cells[1].text = "Groq LLM Reasoning & Fallback"
            table.rows[4].cells[2].text = "Translates regression forecasts and OR-Tools MILP scores into structured decision explanations."
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Gemini API connectors' in p.text:
                    p.text = 'Groq LLM Gateway explaining score increases dynamically with local grounded fallback.'
                if 'Redis cache setups' in p.text:
                    p.text = 'MongoDB native indexes & MinIO S3 proof storage managing fast metric retrieval.'

    # Slide 13: Team Structure
    s13 = prs.slides[12]
    for shape in s13.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Architected Spring Boot gateway services and PostgreSQL' in p.text:
                    p.text = '• Architected Spring Boot 3.2 REST services and MongoDB 2dsphere GeoJSON schemas.'
                elif 'Configured Google OR-Tools solvers and Postgres/PostGIS' in p.text:
                    p.text = '• Configured Google OR-Tools MILP solver, FastAPI XGBoost pipeline, and Groq AI LLM Gateway.'
                elif 'Designed the layout, user interactions' in p.text:
                    p.text = '• Designed Emerald Oasis UI/UX theme, React 18 2D Spatial Canvas, and Groq AI Omnibar.'
                elif 'Built React.js client screens' in p.text:
                    p.text = '• Built 12 workspace nodes, 3-mode view switcher (Canvas, Desktop, Mobile), and interactive charts.'
                elif 'Wrote unit tests and executed integration' in p.text:
                    p.text = '• Built JUnit 5 REST suite, Pytest ML suite, and 36-endpoint Node.js integration test suite.'
                elif 'Created automation scripts for test suites' in p.text:
                    p.text = '• Implemented Server-Sent Events (SSE) notification streaming and system telemetry monitoring.'
                elif 'Supported AI/ML pipeline deployments' in p.text:
                    p.text = '• Configured MinIO S3 object storage and Tesseract OCR utility bill parsing.'
                elif 'Managed frontend-backend integration' in p.text:
                    p.text = '• Built PowerShell multi-service system launcher (run-phase11-system.ps1) and PDF generator.'

    # Slide 14: Implementation Roadmap
    s14 = prs.slides[13]
    for shape in s14.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'Prototype Roadmap' in p.text:
                    p.text = 'Full System Execution Roadmap — ALL PHASES COMPLETED ✅'
                elif 'The implementation roadmap maps 6 weeks' in p.text:
                    p.text = 'The system development roadmap maps 7 complete phases from ideation to production verification.'
                elif 'Phase 1: DB & API Skeleton' in p.text:
                    p.text = 'Phase 1 & 2: Client Spatial UI & Spring Boot Backend — COMPLETED ✅'
                elif 'Initialize Spring Boot API gateway' in p.text:
                    p.text = 'React 18 Spatial Canvas shell, Emerald Oasis theme, Spring Boot 3.2 REST services, JWT auth, and MongoDB schemas.'
                elif 'Phase 2: ML & Data' in p.text:
                    p.text = 'Phase 3 & 4: ML Forecasting, MILP Solver & GIS Tier — COMPLETED ✅'
                elif 'Integrate Python FastAPI ML container' in p.text:
                    p.text = 'FastAPI XGBoost 30-day forecaster, IsolationForest anomaly engine, Google OR-Tools MILP solver, MongoDB 2dsphere geofencing, MinIO S3 OCR.'
                elif 'Phase 3: Digital Twin & Optimization' in p.text:
                    p.text = 'Phase 5 & 6: Groq LLM Gateway, SSE Stream & Telemetry — COMPLETED ✅'
                elif 'Assemble the Household Digital Twin' in p.text:
                    p.text = 'Groq AI Gateway (Llama-3.3-70b) with local grounded fallback engine, SSE real-time notifications, system health telemetry.'
                elif 'Phase 4: AI Assistant & Launch' in p.text:
                    p.text = 'Phase 7: Automated Verification & Production Release — COMPLETED ✅'
                elif 'Connect Gemini API to explain scores' in p.text:
                    p.text = 'JUnit 5 + Pytest + 36-endpoint integration test suite, PowerShell system launcher script (run-phase11-system.ps1).'

    # Slide 15: R&D Highlights
    s15 = prs.slides[14]
    for shape in s15.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'PostgreSQL' in p.text:
                    p.text = p.text.replace('PostgreSQL', 'MongoDB 2dsphere')
                if 'Gemini' in p.text:
                    p.text = p.text.replace('Gemini', 'Groq LLM Gateway')

    # Slide 16: Panel Assessment
    s16 = prs.slides[15]
    for shape in s16.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'PROTOTYPING & R&D' in p.text:
                    p.text = 'PRODUCTION RELEASE & DEPLOYMENT'
                elif 'The EcoSphere Intelligent Environmental Decision Support proposal is highly viable' in p.text:
                    p.text = 'The EcoSphere (VerdantIQ) Intelligent Environmental Decision Support Platform has been fully developed, integrated, and verified across all 36 microservice endpoints. Built on a modular, enterprise stack (React 18 Spatial Canvas, Spring Boot 3.2, MongoDB 2dsphere, FastAPI, Google OR-Tools MILP, Groq AI LLM Gateway, and MinIO S3 OCR), it successfully delivers behavior learning, time-series forecasting, multi-resource solving, and verified activity loops.'
                elif '✓  Proven PostgreSQL relational schema' in p.text:
                    p.text = '✓  Proven MongoDB 2dsphere GeoJSON spatial database'
                elif '✓  Gemini explanation generation prompts' in p.text:
                    p.text = '✓  Groq LLM Gateway & local grounded fallback engine'

    # Slide 17: Project Sanctioned / Final Slide
    s17 = prs.slides[16]
    for shape in s17.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'EcoSphere moves to Phase 1' in p.text:
                    p.text = 'EcoSphere (VerdantIQ) System — Fully Completed & Verified ✅'
                elif 'Focus: DB Design & Core API Gateways' in p.text:
                    p.text = 'Status: 100% Endpoint Verification | Stack: React Spatial UI, Spring Boot 3.2, FastAPI, MongoDB, MinIO S3, Groq AI | Automated PowerShell System Launcher Ready'

    prs.save(pptx_path)
    print(f"Successfully updated PPTX file: {pptx_path}")


if __name__ == '__main__':
    update_docx()
    update_pptx()
